"""
応募→面接→研修→配布員登録 業務フロー＆マニュアル PPTX 生成スクリプト
LINE連携・Posting System同期を含む
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Colors ──
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1E, 0x29, 0x3B)
DARK_GRAY = RGBColor(0x33, 0x44, 0x55)
MEDIUM_GRAY = RGBColor(0x64, 0x74, 0x8B)
LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
ACCENT_BLUE = RGBColor(0x37, 0x56, 0xE8)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
ACCENT_AMBER = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)
ACCENT_INDIGO = RGBColor(0x63, 0x66, 0xF1)
ACCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
LINE_GREEN = RGBColor(0x06, 0xC7, 0x55)
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
BG_BLUE = RGBColor(0x1E, 0x3A, 0x5F)

def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14, color=BLACK, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox

def add_rich_text_box(slide, left, top, width, height, lines):
    """lines = [(text, font_size, color, bold, alignment), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, font_size, color, bold, alignment) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = alignment
        p.space_after = Pt(4)
    return txBox

def add_arrow(slide, left, top, width, height, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_down_arrow(slide, left, top, width, height, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def make_step_card(slide, left, top, width, height, icon, title, desc, color):
    """Create a step card with icon, title, and description"""
    card = add_shape(slide, left, top, width, height, WHITE, color, Pt(2))
    # Color bar at top
    bar = add_shape(slide, left + Inches(0.05), top + Inches(0.05), width - Inches(0.1), Inches(0.35), color)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.07), width - Inches(0.3), Inches(0.35),
                 f"{icon}  {title}", 13, WHITE, True, PP_ALIGN.CENTER)
    # Description
    add_text_box(slide, left + Inches(0.15), top + Inches(0.5), width - Inches(0.3), height - Inches(0.6),
                 desc, 10, DARK_GRAY, False, PP_ALIGN.LEFT)
    return card

# ═══════════════════════════════════════════════════════════════
# SLIDE 1: Cover
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = BG_DARK

# Accent bar
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT_BLUE)

# Title block
add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.0),
             "採用業務フロー＆オペレーションマニュアル", 38, WHITE, True, PP_ALIGN.LEFT)
add_text_box(slide, Inches(1.5), Inches(2.9), Inches(10), Inches(0.6),
             "Recruitment Workflow & Operations Manual", 20, MEDIUM_GRAY, False, PP_ALIGN.LEFT)

# Divider
add_shape(slide, Inches(1.5), Inches(3.8), Inches(3), Inches(0.04), ACCENT_BLUE)

add_text_box(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(0.5),
             "応募受付 → 面接 → 採用 → 研修 → 配布員登録 → LINE連携 → アプリ配信", 16, MEDIUM_GRAY, False, PP_ALIGN.LEFT)

# Footer
add_text_box(slide, Inches(1.5), Inches(5.5), Inches(5), Inches(0.4),
             "株式会社ティラミス / K&Partners", 14, MEDIUM_GRAY, False, PP_ALIGN.LEFT)
add_text_box(slide, Inches(1.5), Inches(5.9), Inches(5), Inches(0.4),
             "PMS Pro — ポスティング管理システム", 12, MEDIUM_GRAY, False, PP_ALIGN.LEFT)
add_text_box(slide, Inches(1.5), Inches(6.3), Inches(5), Inches(0.4),
             "2026年3月", 12, MEDIUM_GRAY, False, PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2: Full Flow Overview
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = WHITE

add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), BG_DARK)
add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
             "全体業務フロー  Overall Recruitment Workflow", 24, WHITE, True, PP_ALIGN.LEFT)

# Step boxes in a flow
steps = [
    ("1", "応募受付", "Apply", ACCENT_BLUE),
    ("2", "面接", "Interview", ACCENT_INDIGO),
    ("3", "採用判定", "Hire", ACCENT_GREEN),
    ("4", "研修", "Training", ACCENT_AMBER),
    ("5", "配布員登録", "Register", ACCENT_RED),
    ("6", "LINE連携", "LINE", LINE_GREEN),
    ("7", "アプリ配信", "App", ACCENT_PURPLE),
]

start_x = Inches(0.5)
y = Inches(1.8)
box_w = Inches(1.5)
box_h = Inches(1.5)
gap = Inches(0.25)
arrow_w = Inches(0.3)

for i, (num, label_ja, label_en, color) in enumerate(steps):
    x = start_x + i * (box_w + gap + arrow_w)
    card = add_shape(slide, x, y, box_w, box_h, color)
    # Number circle
    add_text_box(slide, x, y + Inches(0.15), box_w, Inches(0.5),
                 f"STEP {num}", 12, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(0.5), box_w, Inches(0.5),
                 label_ja, 18, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(0.95), box_w, Inches(0.4),
                 label_en, 11, RGBColor(0xDD, 0xDD, 0xFF), False, PP_ALIGN.CENTER)

    # Arrow between steps
    if i < len(steps) - 1:
        ax = x + box_w + Inches(0.02)
        add_arrow(slide, ax, y + Inches(0.5), arrow_w - Inches(0.04), Inches(0.4), MEDIUM_GRAY)

# Detail rows below
detail_y = Inches(3.7)
details = [
    ("応募者が /apply から応募\n面接スロット自動選択\nGoogle Meet自動生成", ACCENT_BLUE),
    ("管理者が /applicants で\n面接実施・評価入力\n日程変更/キャンセル可", ACCENT_INDIGO),
    ("採用/不採用を判定\n採用→研修ステップへ\n不採用→メール通知", ACCENT_GREEN),
    ("研修スロット割り当て\n または自己予約リンク送信\nグループ研修（定員制）", ACCENT_AMBER),
    ("PMS配布員として登録\nPosting Systemに同期\n契約書自動送信", ACCENT_RED),
    ("LINE公式アカウント\nフォロワー取込＆紐付け\n通知基盤構築", LINE_GREEN),
    ("TestFlight招待(iOS)\nGoogle Play案内(Android)\nポータル情報送信", ACCENT_PURPLE),
]

for i, (detail, color) in enumerate(details):
    x = start_x + i * (box_w + gap + arrow_w)
    # Thin color line
    add_shape(slide, x, detail_y, box_w, Inches(0.04), color)
    add_text_box(slide, x, detail_y + Inches(0.15), box_w, Inches(1.5),
                 detail, 9, DARK_GRAY, False, PP_ALIGN.LEFT)

# Systems involved
sys_y = Inches(5.5)
add_shape(slide, Inches(0.3), sys_y, Inches(12.7), Inches(1.7), LIGHT_GRAY)
add_text_box(slide, Inches(0.5), sys_y + Inches(0.1), Inches(12), Inches(0.4),
             "関連システム  Related Systems", 14, DARK_GRAY, True, PP_ALIGN.LEFT)

systems = [
    ("PMS Pro", "社内基幹管理システム\npms.tiramis.co.jp", ACCENT_BLUE),
    ("Posting System", "配布管理（旧システム）\nm_staff テーブル同期", ACCENT_RED),
    ("LINE Official", "K&Partners公式アカウント\n@309nmzqx", LINE_GREEN),
    ("Google Calendar", "Meet自動生成\n面接リンク作成", RGBColor(0x42, 0x85, 0xF4)),
    ("DocuSeal", "業務委託契約書\n電子署名", ACCENT_PURPLE),
    ("App Store / Play", "TestFlight招待\nGoogle Play案内", DARK_GRAY),
]

for i, (name, desc, color) in enumerate(systems):
    sx = Inches(0.5) + i * Inches(2.1)
    add_shape(slide, sx, sys_y + Inches(0.55), Inches(1.9), Inches(0.04), color)
    add_text_box(slide, sx, sys_y + Inches(0.65), Inches(1.9), Inches(0.3),
                 name, 11, color, True, PP_ALIGN.LEFT)
    add_text_box(slide, sx, sys_y + Inches(0.95), Inches(1.9), Inches(0.6),
                 desc, 9, MEDIUM_GRAY, False, PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3: Step 1 - Application
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = WHITE

add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), ACCENT_BLUE)
add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
             "STEP 1  応募受付  Application", 24, WHITE, True, PP_ALIGN.LEFT)

# Left: Flow
add_shape(slide, Inches(0.3), Inches(1.2), Inches(6.2), Inches(6.0), LIGHT_GRAY)
add_text_box(slide, Inches(0.5), Inches(1.3), Inches(5.5), Inches(0.4),
             "応募フロー", 16, DARK_GRAY, True, PP_ALIGN.LEFT)

flow_items = [
    ("1", "応募者が /apply にアクセス", "求人媒体のURLに ?source=xxx パラメータを付けて\n応募経路を自動追跡（例: ?source=indeed）"),
    ("2", "応募フォームに入力", "氏名、電話番号、メール、国籍、ビザ種類\n希望職種、面接希望日時を選択"),
    ("3", "面接スロットを選択", "職種に対応する空きスロットのみ表示\n全職種対応スロット＋職種別スロット"),
    ("4", "送信", "Google Meet自動生成（環境変数設定時）\n確認メールを応募者に送信\nスロットが「予約済み」に更新"),
]

for i, (num, title, desc) in enumerate(flow_items):
    fy = Inches(1.9) + i * Inches(1.2)
    add_shape(slide, Inches(0.5), fy, Inches(0.4), Inches(0.4), ACCENT_BLUE)
    add_text_box(slide, Inches(0.5), fy + Inches(0.02), Inches(0.4), Inches(0.4),
                 num, 14, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.05), fy, Inches(5.2), Inches(0.35),
                 title, 13, BLACK, True, PP_ALIGN.LEFT)
    add_text_box(slide, Inches(1.05), fy + Inches(0.35), Inches(5.2), Inches(0.7),
                 desc, 10, MEDIUM_GRAY, False, PP_ALIGN.LEFT)

# Right: Admin Setup
add_shape(slide, Inches(6.8), Inches(1.2), Inches(6.2), Inches(6.0), LIGHT_GRAY)
add_text_box(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.4),
             "管理者の事前設定", 16, DARK_GRAY, True, PP_ALIGN.LEFT)

setup_items = [
    ("面接スロット設定", "/settings → 面接スロットタブ\n・デフォルトスロットマスタ（曜日×時間帯）を設定\n・CRONで14日分を自動生成（毎日01:00）\n・個別スロットの手動追加も可能\n・職種別 or 全職種対応を選択"),
    ("職種マスタ設定", "/settings → 職種タブ\n・職種名（日/英）を登録\n・応募フォームの選択肢に反映"),
    ("求人媒体設定", "/settings → 求人媒体タブ\n・媒体名とコード（URLパラメータ値）を登録\n・例: code=indeed, nameJa=Indeed\n・自動的に小文字に正規化"),
    ("応募フォームURL", "基本URL: https://pms.tiramis.co.jp/apply\n求人媒体追跡: ?source=indeed\n言語対応: 日本語/英語自動切替"),
]

for i, (title, desc) in enumerate(setup_items):
    sy = Inches(1.9) + i * Inches(1.35)
    add_text_box(slide, Inches(7.0), sy, Inches(5.8), Inches(0.3),
                 f"■ {title}", 12, ACCENT_BLUE, True, PP_ALIGN.LEFT)
    add_text_box(slide, Inches(7.2), sy + Inches(0.3), Inches(5.6), Inches(1.0),
                 desc, 10, MEDIUM_GRAY, False, PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4: Step 2 - Interview
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = WHITE

add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), ACCENT_INDIGO)
add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
             "STEP 2  面接・評価  Interview & Evaluation", 24, WHITE, True, PP_ALIGN.LEFT)

# Main content area
add_text_box(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(0.4),
             "面接管理画面  /applicants", 16, DARK_GRAY, True, PP_ALIGN.LEFT)

# Left column: Calendar + List
left_items = [
    ("FullCalendar カレンダー", "・月/週/日ビューで面接スケジュールを可視化\n・スロット上に応募者名を表示\n・クリックで詳細モーダルを開く"),
    ("応募者一覧テーブル", "・ステータスフィルタ: 面接待ち / 研修待ち / 研修完了\n・採用状態フィルタ: 選考中 / 採用 / 不採用\n・検索: 名前、メール、電話番号"),
    ("評価モーダル", "・個人情報（氏名、連絡先、国籍、ビザ）\n・面接日時 + Google Meet リンク\n・応募経路（求人媒体）表示・変更\n・評価スコア入力（1-5段階、カラーコード）\n・メモ / コメント入力"),
]

for i, (title, desc) in enumerate(left_items):
    ly = Inches(1.8) + i * Inches(1.7)
    add_shape(slide, Inches(0.5), ly, Inches(6), Inches(1.5), LIGHT_GRAY, ACCENT_INDIGO, Pt(1))
    add_text_box(slide, Inches(0.7), ly + Inches(0.1), Inches(5.5), Inches(0.3),
                 title, 13, ACCENT_INDIGO, True, PP_ALIGN.LEFT)
    add_text_box(slide, Inches(0.7), ly + Inches(0.4), Inches(5.5), Inches(1.0),
                 desc, 10, DARK_GRAY, False, PP_ALIGN.LEFT)

# Right column: Actions
add_text_box(slide, Inches(7.0), Inches(1.2), Inches(6), Inches(0.4),
             "管理者アクション", 16, DARK_GRAY, True, PP_ALIGN.LEFT)

right_items = [
    ("面接日程変更", "・評価モーダル → 「日程変更」ボタン\n・空きスロット一覧から新日程を選択\n・旧スロット自動解放 + 新スロット予約\n・Google Meet自動再生成", ACCENT_BLUE),
    ("面接キャンセル", "・評価モーダル → 「面接キャンセル」ボタン\n・確認ダイアログ後に実行\n・スロットを解放（他の応募者が使える）\n・ステータスは変更しない", ACCENT_AMBER),
    ("採用/不採用判定", "・評価モーダル → 採用ステータス変更\n  IN_PROGRESS → HIRED or REJECTED\n・採用通知メール送信\n・採用 → Step 3 (研修) へ進む", ACCENT_GREEN),
    ("ステータス管理", "・INTERVIEW_WAITING: 面接待ち\n・TRAINING_WAITING: 研修待ち\n・TRAINING_COMPLETED: 研修完了\n\n採用状態:\n・IN_PROGRESS → HIRED / REJECTED", MEDIUM_GRAY),
]

for i, (title, desc, color) in enumerate(right_items):
    ry = Inches(1.8) + i * Inches(1.35)
    add_shape(slide, Inches(7.0), ry, Inches(6), Inches(1.2), WHITE, color, Pt(2))
    add_text_box(slide, Inches(7.2), ry + Inches(0.05), Inches(5.5), Inches(0.3),
                 title, 12, color, True, PP_ALIGN.LEFT)
    add_text_box(slide, Inches(7.2), ry + Inches(0.3), Inches(5.5), Inches(0.85),
                 desc, 9.5, DARK_GRAY, False, PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5: Step 3 & 4 - Training
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = WHITE

add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), ACCENT_AMBER)
add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
             "STEP 3-4  採用→研修  Hiring → Training", 24, WHITE, True, PP_ALIGN.LEFT)

# Left: Training slot management
add_shape(slide, Inches(0.3), Inches(1.2), Inches(6.2), Inches(6.0), LIGHT_GRAY)
add_text_box(slide, Inches(0.5), Inches(1.3), Inches(5.5), Inches(0.4),
             "研修スロット管理", 16, DARK_GRAY, True, PP_ALIGN.LEFT)

training_items = [
    ("研修スロットマスタ", "/settings → 研修スロットタブ\n・曜日ごとの開始/終了時刻、間隔、定員を設定\n・CRONで14日分を自動生成（毎日02:00）\n・有効/無効の切り替え可能"),
    ("研修スロット一覧", "/training-slots\n・手動でスロットを追加可能\n・応募者紐付き時は削除不可\n・remainingCapacity（残席）を表示"),
    ("研修割り当て方法①", "管理者が直接割り当て:\n評価モーダル → 「研修スロット」セクション\n→ 空きスロット一覧から選択\n→ 確認メール自動送信"),
    ("研修割り当て方法②", "応募者に自己予約させる:\n評価モーダル → 「後でメール案内」\n→ 自己予約リンク付きメール送信\n→ 応募者が /training-booking?token=xxx\n  で自分で予約（日/英対応）"),
]

for i, (title, desc) in enumerate(training_items):
    ty = Inches(1.9) + i * Inches(1.35)
    add_text_box(slide, Inches(0.5), ty, Inches(5.8), Inches(0.3),
                 f"■ {title}", 12, ACCENT_AMBER, True, PP_ALIGN.LEFT)
    add_text_box(slide, Inches(0.7), ty + Inches(0.3), Inches(5.6), Inches(1.0),
                 desc, 10, MEDIUM_GRAY, False, PP_ALIGN.LEFT)

# Right: Flow diagram
add_shape(slide, Inches(6.8), Inches(1.2), Inches(6.2), Inches(6.0), LIGHT_GRAY)
add_text_box(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.4),
             "採用後の流れ", 16, DARK_GRAY, True, PP_ALIGN.LEFT)

flow_steps = [
    ("採用決定", "hiringStatus → HIRED", ACCENT_GREEN),
    ("研修スロット割り当て", "管理者直接 or 自己予約リンク", ACCENT_AMBER),
    ("研修実施", "flowStatus → TRAINING_COMPLETED", ACCENT_BLUE),
    ("配布員登録可能に", "「配布員登録」ボタンが表示される\n条件: HIRED + 研修スロット設定済み", ACCENT_RED),
]

for i, (title, desc, color) in enumerate(flow_steps):
    fy = Inches(1.9) + i * Inches(1.3)
    add_shape(slide, Inches(7.2), fy, Inches(5.5), Inches(0.9), WHITE, color, Pt(2))
    add_text_box(slide, Inches(7.4), fy + Inches(0.05), Inches(5), Inches(0.3),
                 f"Step {i+1}: {title}", 12, color, True, PP_ALIGN.LEFT)
    add_text_box(slide, Inches(7.4), fy + Inches(0.35), Inches(5), Inches(0.5),
                 desc, 10, DARK_GRAY, False, PP_ALIGN.LEFT)
    if i < len(flow_steps) - 1:
        add_down_arrow(slide, Inches(9.7), fy + Inches(0.95), Inches(0.3), Inches(0.25), MEDIUM_GRAY)

# Note
add_shape(slide, Inches(7.0), Inches(6.0), Inches(5.8), Inches(1.0), RGBColor(0xFF, 0xFB, 0xEB), ACCENT_AMBER, Pt(1))
add_text_box(slide, Inches(7.2), Inches(6.05), Inches(5.4), Inches(0.9),
             "注意: 研修スロットが未設定の場合、\n「配布員登録」ボタンは表示されません。\n必ず研修スロットを割り当ててから登録に進んでください。", 10, DARK_GRAY, False, PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6: Step 5 - Distributor Registration + Posting System
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = WHITE

add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), ACCENT_RED)
add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
             "STEP 5  配布員登録 & Posting System同期  Distributor Registration", 24, WHITE, True, PP_ALIGN.LEFT)

# Left: Registration form
add_shape(slide, Inches(0.3), Inches(1.2), Inches(6.2), Inches(6.0), LIGHT_GRAY)
add_text_box(slide, Inches(0.5), Inches(1.3), Inches(5.5), Inches(0.4),
             "配布員登録フォーム", 16, DARK_GRAY, True, PP_ALIGN.LEFT)

add_text_box(slide, Inches(0.5), Inches(1.8), Inches(5.8), Inches(0.3),
             "評価モーダル → セクション6「配布員登録」", 11, ACCENT_RED, True, PP_ALIGN.LEFT)

reg_items = [
    ("表示条件", "hiringStatus === 'HIRED'\n&& applicant.trainingSlot !== null"),
    ("必須入力", "・生年月日（初期PW = YYYYMMDD のハッシュ）\n・所属支店（ドロップダウン）"),
    ("任意入力", "・社員番号（staffId）— 未入力時は自動採番\n・性別"),
    ("自動継承", "応募者情報から自動コピー:\n氏名、メール、電話、国籍、ビザ種類、\n住所、郵便番号、ビル名"),
    ("Posting System同期", "チェックボックス「Posting Systemにも登録する」\nON → 登録と同時にPosting Systemのm_staffに同期\nOFF → PMSのみ登録"),
    ("DocuSeal契約書", "チェックボックス「業務委託契約書を自動送信する」\nON → DocuSeal経由で電子署名依頼を自動送信\n配布員のメールアドレスに契約書が届く"),
]

for i, (title, desc) in enumerate(reg_items):
    ry = Inches(2.2) + i * Inches(0.9)
    add_text_box(slide, Inches(0.5), ry, Inches(5.8), Inches(0.25),
                 f"● {title}", 11, ACCENT_RED, True, PP_ALIGN.LEFT)
    add_text_box(slide, Inches(0.8), ry + Inches(0.25), Inches(5.5), Inches(0.65),
                 desc, 9.5, DARK_GRAY, False, PP_ALIGN.LEFT)

# Right: Posting System Sync details
add_shape(slide, Inches(6.8), Inches(1.2), Inches(6.2), Inches(3.5), WHITE, ACCENT_RED, Pt(2))
add_text_box(slide, Inches(7.0), Inches(1.3), Inches(5.8), Inches(0.4),
             "Posting System 同期の仕組み", 16, ACCENT_RED, True, PP_ALIGN.LEFT)

sync_desc = """■ 同期先
Posting System の m_staff テーブル
(postingrealdb.m_staff)

■ 同期データ
・STAFF_CD: 社員番号（7文字以内）
・STAFF_NAME: 氏名
・STAFF_TEL: 電話番号
・SHOP_CD: 支店コード（自動変換）
・JOIN_DATE: 入社日
・STAFF_DUTY_DIV: 1（ポスティングスタッフ）

■ 支店コード変換
高田馬場→MBF  横浜→MYO  新松戸→MMA
浦和→MUr  西新井→MNA  新小岩→MKO
蒲田→MKA  赤羽→MAK  吉祥寺→MKI

■ 動作
INSERT（新規）or UPDATE STAFF_TEL（既存）
fire-and-forget方式（失敗してもPMS登録は完了）"""

add_text_box(slide, Inches(7.0), Inches(1.8), Inches(5.8), Inches(2.7),
             sync_desc, 9.5, DARK_GRAY, False, PP_ALIGN.LEFT)

# Bottom right: After registration
add_shape(slide, Inches(6.8), Inches(5.0), Inches(6.2), Inches(2.2), LIGHT_GRAY)
add_text_box(slide, Inches(7.0), Inches(5.1), Inches(5.5), Inches(0.3),
             "登録完了後", 14, DARK_GRAY, True, PP_ALIGN.LEFT)

after_items = [
    "・配布員IDが評価モーダルに表示される",
    "・配布員は /staff/login からログイン可能",
    "・初期パスワード: 生年月日（YYYYMMDD）",
    "・初回ログイン時にパスワード変更を強制",
    "・Posting System側でもスケジュール割当が可能に",
]

for i, item in enumerate(after_items):
    add_text_box(slide, Inches(7.0), Inches(5.5) + i * Inches(0.3), Inches(5.8), Inches(0.3),
                 item, 10, DARK_GRAY, False, PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7: Step 6 - LINE Integration
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = WHITE

add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), LINE_GREEN)
add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
             "STEP 6  LINE連携管理  LINE Integration", 24, WHITE, True, PP_ALIGN.LEFT)

# Left: Overview
add_shape(slide, Inches(0.3), Inches(1.2), Inches(6.2), Inches(3.0), LIGHT_GRAY)
add_text_box(slide, Inches(0.5), Inches(1.3), Inches(5.5), Inches(0.4),
             "LINE連携の概要", 16, DARK_GRAY, True, PP_ALIGN.LEFT)

line_overview = """・LINE公式アカウント「K&Partners」(@309nmzqx) を使用
・配布員のLINEアカウントとPMS上の配布員情報を紐付け
・将来的にスケジュール通知やお知らせをLINEで配信

■ 仕組み
1. 配布員がLINE公式アカウントを友だち追加
2. 自動で登録依頼メッセージ（日英バイリンガル）が送信
3. ボタンをタップ → Webhook経由でPMSに登録
4. 管理者が /line 画面で配布員と紐付け"""

add_text_box(slide, Inches(0.5), Inches(1.8), Inches(5.8), Inches(2.2),
             line_overview, 10, DARK_GRAY, False, PP_ALIGN.LEFT)

# Right top: Management UI
add_shape(slide, Inches(6.8), Inches(1.2), Inches(6.2), Inches(3.0), LIGHT_GRAY)
add_text_box(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.4),
             "LINE連携管理画面  /line", 16, DARK_GRAY, True, PP_ALIGN.LEFT)

mgmt_desc = """■ 画面構成
・統計カード: 合計 / 紐付け済 / 未紐付け
・フィルタ: すべて / 紐付け済み / 未紐付け
・検索: LINE表示名で絞り込み
・一覧テーブル: LINEアイコン＋表示名、フォロー状態、紐付け配布員

■ 操作
・「LINE登録依頼を送信」ボタン
  → 全フォロワーにFlexメッセージを一斉送信
・「紐付け」ボタン → 配布員を検索して選択
・「解除」ボタン → 確認後に紐付け解除"""

add_text_box(slide, Inches(7.0), Inches(1.8), Inches(5.8), Inches(2.2),
             mgmt_desc, 10, DARK_GRAY, False, PP_ALIGN.LEFT)

# Bottom: Webhook flow
add_shape(slide, Inches(0.3), Inches(4.5), Inches(12.7), Inches(2.8), WHITE, LINE_GREEN, Pt(2))
add_text_box(slide, Inches(0.5), Inches(4.6), Inches(12), Inches(0.4),
             "Webhook 自動処理フロー", 14, LINE_GREEN, True, PP_ALIGN.LEFT)

webhook_steps = [
    ("友だち追加\n(follow)", "ユーザー登録\n+ 連携依頼送信", "ボタンタップ\n(postback)", "PMS DBに\nユーザー保存", "管理者が\n配布員と紐付け"),
]

labels = ["友だち追加\n(follow)", "プロフィール取得\n& DB保存", "連携依頼\nメッセージ送信", "ボタンタップ\n(postback)", "管理者が\n紐付け実行"]
colors_wh = [LINE_GREEN, ACCENT_BLUE, ACCENT_INDIGO, ACCENT_GREEN, ACCENT_RED]

for i, (label, color) in enumerate(zip(labels, colors_wh)):
    wx = Inches(0.5) + i * Inches(2.5)
    add_shape(slide, wx, Inches(5.2), Inches(2.0), Inches(1.0), color)
    add_text_box(slide, wx, Inches(5.3), Inches(2.0), Inches(0.8),
                 label, 10, WHITE, True, PP_ALIGN.CENTER)
    if i < len(labels) - 1:
        add_arrow(slide, wx + Inches(2.05), Inches(5.5), Inches(0.35), Inches(0.3), MEDIUM_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8: Step 7 - App Distribution + Portal Notification
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = WHITE

add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), ACCENT_PURPLE)
add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
             "STEP 7  アプリ配信 & ポータル通知  App Distribution", 24, WHITE, True, PP_ALIGN.LEFT)

# Left: App Distribution
add_shape(slide, Inches(0.3), Inches(1.2), Inches(6.2), Inches(5.8), LIGHT_GRAY)
add_text_box(slide, Inches(0.5), Inches(1.3), Inches(5.5), Inches(0.4),
             "アプリ配信  /distributors/[id]", 16, DARK_GRAY, True, PP_ALIGN.LEFT)

app_items = [
    ("アプリ配信ボタン", "配布員詳細画面のヘッダーに「アプリ配信」ボタン\nクリックで配信モーダルが開く"),
    ("プラットフォーム選択", "・Apple (TestFlight)\n  → App Store Connect API で TestFlight 招待送信\n  → 既に登録済みの場合は自動で再招待\n\n・Android (Google Play)\n  → インストール案内メールを送信"),
    ("メールアドレス", "配布員のメールアドレスが初期値\n別のアドレスに変更も可能"),
    ("配信履歴", "モーダル下部に過去の配信履歴を表示\nステータス: PENDING / SENT / FAILED"),
]

for i, (title, desc) in enumerate(app_items):
    ay = Inches(1.9) + i * Inches(1.3)
    add_text_box(slide, Inches(0.5), ay, Inches(5.8), Inches(0.25),
                 f"● {title}", 11, ACCENT_PURPLE, True, PP_ALIGN.LEFT)
    add_text_box(slide, Inches(0.8), ay + Inches(0.25), Inches(5.5), Inches(1.0),
                 desc, 9.5, DARK_GRAY, False, PP_ALIGN.LEFT)

# Right: Portal Notification + Summary
add_shape(slide, Inches(6.8), Inches(1.2), Inches(6.2), Inches(2.5), LIGHT_GRAY)
add_text_box(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.4),
             "ポータル通知メール", 16, DARK_GRAY, True, PP_ALIGN.LEFT)

portal_desc = """配布員詳細画面の「ポータル通知」ボタンで送信

■ 送信内容
・ログインURL: https://pms.tiramis.co.jp/staff/login
・社員番号（ログインID）
・初期パスワード: 生年月日（YYYYMMDD）
・初回ログイン時にパスワード変更が必要な旨

■ パスワードリセット
・管理者が「PW」ボタンからリセット可能
・生年月日（YYYYMMDD）にリセットされる"""

add_text_box(slide, Inches(7.0), Inches(1.8), Inches(5.8), Inches(1.8),
             portal_desc, 10, DARK_GRAY, False, PP_ALIGN.LEFT)

# Checklist
add_shape(slide, Inches(6.8), Inches(4.0), Inches(6.2), Inches(3.2), WHITE, ACCENT_GREEN, Pt(2))
add_text_box(slide, Inches(7.0), Inches(4.1), Inches(5.5), Inches(0.4),
             "配布員セットアップ完了チェックリスト", 14, ACCENT_GREEN, True, PP_ALIGN.LEFT)

checklist = [
    "□ 配布員としてPMSに登録済み",
    "□ Posting Systemに同期済み（チェックONの場合）",
    "□ DocuSeal業務委託契約書を送信済み",
    "□ LINE公式アカウントと紐付け済み",
    "□ スタッフポータル案内メールを送信済み",
    "□ アプリ招待を送信済み（iOS / Android）",
    "□ 初回ログイン確認（パスワード変更済み）",
]

for i, item in enumerate(checklist):
    add_text_box(slide, Inches(7.2), Inches(4.6) + i * Inches(0.33), Inches(5.5), Inches(0.3),
                 item, 11, DARK_GRAY, False, PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════
output_dir = os.path.expanduser("~/Downloads/PMS_設計書")
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, "PMS_採用業務フロー_マニュアル.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
