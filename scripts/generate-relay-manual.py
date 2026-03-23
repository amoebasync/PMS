#!/usr/bin/env python3
"""中継/回収 担当者向け業務マニュアル PowerPoint 生成スクリプト"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# === Color Palette ===
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1E, 0x29, 0x3B)
DARK_GRAY = RGBColor(0x47, 0x55, 0x69)
MID_GRAY = RGBColor(0x94, 0xA3, 0xB8)
LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
VERY_LIGHT = RGBColor(0xF8, 0xFA, 0xFC)
INDIGO = RGBColor(0x63, 0x66, 0xF1)
INDIGO_DARK = RGBColor(0x49, 0x46, 0xCD)
INDIGO_LIGHT = RGBColor(0xE0, 0xE7, 0xFF)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
ORANGE_LIGHT = RGBColor(0xFF, 0xED, 0xD5)
PURPLE = RGBColor(0x9F, 0x33, 0xE2)  # reduced blue for proper purple
PURPLE_LIGHT = RGBColor(0xF3, 0xE8, 0xFF)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
GREEN_LIGHT = RGBColor(0xDC, 0xFC, 0xE7)
RED = RGBColor(0xEF, 0x44, 0x44)
RED_LIGHT = RGBColor(0xFE, 0xE2, 0xE2)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
AMBER_LIGHT = RGBColor(0xFE, 0xF3, 0xC7)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
BLUE_LIGHT = RGBColor(0xDB, 0xEA, 0xFE)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(0), corner_radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=14, color=BLACK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Yu Gothic UI'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_badge(slide, left, top, text, bg_color, text_color, width=Inches(1.4), height=Inches(0.4), font_size=11):
    shape = add_rect(slide, left, top, width, height, fill_color=bg_color, corner_radius=True)
    shape.text_frame.word_wrap = True
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = 'Yu Gothic UI'
    p.alignment = PP_ALIGN.CENTER
    shape.text_frame.paragraphs[0].space_before = Pt(0)
    shape.text_frame.paragraphs[0].space_after = Pt(0)
    return shape


def add_icon_circle(slide, left, top, size, color, text, font_size=18):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Yu Gothic UI'
    p.alignment = PP_ALIGN.CENTER
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_arrow_right(slide, left, top, width=Inches(0.6), height=Inches(0.3), color=MID_GRAY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_multiline_text(slide, left, top, width, height, lines, font_size=12, color=DARK_GRAY, line_spacing=1.5, font_name='Yu Gothic UI'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(font_size * (line_spacing - 1))
    return txBox


# =====================================================
# Slide 1: Title
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, VERY_LIGHT)

# Top accent bar
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), fill_color=INDIGO)

# Main title card
card_w = Inches(9)
card_h = Inches(3.8)
card_left = (W - card_w) // 2
card_top = Inches(1.5)
add_rect(slide, card_left, card_top, card_w, card_h, fill_color=WHITE, border_color=RGBColor(0xE2, 0xE8, 0xF0), border_width=Pt(1), corner_radius=True)

# Indigo accent on left side of card
add_rect(slide, card_left, card_top, Inches(0.08), card_h, fill_color=INDIGO)

add_text(slide, card_left + Inches(0.6), card_top + Inches(0.5), Inches(8), Inches(0.5),
         '業務マニュアル', font_size=16, color=INDIGO, bold=True)

add_text(slide, card_left + Inches(0.6), card_top + Inches(1.0), Inches(8), Inches(1.2),
         '中継/回収 担当者ガイド', font_size=40, color=BLACK, bold=True)

add_text(slide, card_left + Inches(0.6), card_top + Inches(2.3), Inches(8), Inches(0.5),
         'Relay & Collection Operations Manual', font_size=18, color=MID_GRAY)

add_text(slide, card_left + Inches(0.6), card_top + Inches(3.0), Inches(8), Inches(0.5),
         '株式会社ティラミス  |  PMS Pro', font_size=14, color=MID_GRAY)


# Three small badges at bottom
badge_y = Inches(5.8)
add_badge(slide, Inches(4.2), badge_y, '中継', ORANGE_LIGHT, ORANGE, Inches(1.2), Inches(0.35))
add_badge(slide, Inches(5.7), badge_y, '回収', PURPLE_LIGHT, PURPLE, Inches(1.2), Inches(0.35))
add_badge(slide, Inches(7.2), badge_y, '全中継', GREEN_LIGHT, GREEN, Inches(1.2), Inches(0.35))


# =====================================================
# Slide 2: Overview - What are Relay/Collection tasks?
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), fill_color=INDIGO)

# Title
add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         '中継/回収とは？', font_size=30, color=BLACK, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.4),
         '配布作業中に発生するチラシの補充・回収業務です', font_size=14, color=MID_GRAY)

# Three cards
card_data = [
    {
        'icon': '🚚', 'title': '中継（RELAY）',
        'desc': '配布員にチラシを届ける',
        'detail': '配布員が配布エリアでチラシが不足した際に、\n担当者がチラシを指定場所まで届けます。',
        'color': ORANGE, 'bg': ORANGE_LIGHT,
    },
    {
        'icon': '📦', 'title': '回収（COLLECTION）',
        'desc': '配布員からチラシを回収する',
        'detail': '配布終了後や余ったチラシを、\n担当者が指定場所で回収します。',
        'color': PURPLE, 'bg': PURPLE_LIGHT,
    },
    {
        'icon': '🔄', 'title': '全中継（FULL RELAY）',
        'desc': '全チラシを一括で届ける',
        'detail': '配布員の全チラシをまとめて\n指定場所まで届ける業務です。',
        'color': GREEN, 'bg': GREEN_LIGHT,
    },
]

for i, data in enumerate(card_data):
    x = Inches(0.8) + Inches(4.0) * i
    y = Inches(1.8)
    cw = Inches(3.7)
    ch = Inches(4.8)

    # Card background
    add_rect(slide, x, y, cw, ch, fill_color=VERY_LIGHT, border_color=RGBColor(0xE2, 0xE8, 0xF0), border_width=Pt(1), corner_radius=True)

    # Top color stripe
    add_rect(slide, x, y, cw, Inches(0.06), fill_color=data['color'])

    # Icon
    add_text(slide, x + Inches(0.3), y + Inches(0.3), Inches(1), Inches(0.8), data['icon'], font_size=36)

    # Badge
    add_badge(slide, x + Inches(0.3), y + Inches(1.2), data['title'], data['bg'], data['color'], Inches(2.8), Inches(0.4), font_size=13)

    # Subtitle
    add_text(slide, x + Inches(0.3), y + Inches(1.9), cw - Inches(0.6), Inches(0.4),
             data['desc'], font_size=14, color=BLACK, bold=True)

    # Detail
    add_multiline_text(slide, x + Inches(0.3), y + Inches(2.5), cw - Inches(0.6), Inches(2.0),
                       data['detail'].split('\n'), font_size=12, color=DARK_GRAY)


# =====================================================
# Slide 3: Overall Flow
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), fill_color=INDIGO)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         '業務フロー全体像', font_size=30, color=BLACK, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.4),
         '管理者からタスクが割り当てられ、完了まで以下のフローで進行します', font_size=14, color=MID_GRAY)

# Flow steps
steps = [
    {'num': '1', 'title': 'タスク確認', 'desc': 'PMS画面で\n当日のタスクを確認', 'color': INDIGO},
    {'num': '2', 'title': 'ルート確認', 'desc': '最適化ルートで\n巡回順序を確認', 'color': BLUE},
    {'num': '3', 'title': 'ナビ開始', 'desc': 'Google Maps で\nナビゲーション開始', 'color': GREEN},
    {'num': '4', 'title': '中継/回収', 'desc': '現場で配布員と\nチラシの受け渡し', 'color': ORANGE},
    {'num': '5', 'title': 'エビデンス', 'desc': '現場の写真を撮影\nPMSにアップロード', 'color': PURPLE},
    {'num': '6', 'title': '完了報告', 'desc': 'ステータスを\n「完了」に変更', 'color': INDIGO_DARK},
]

step_w = Inches(1.7)
total_flow_w = step_w * 6 + Inches(0.5) * 5
start_x = (W - total_flow_w) // 2
y_top = Inches(2.0)

for i, step in enumerate(steps):
    x = start_x + (step_w + Inches(0.5)) * i

    # Circle with number
    circle_size = Inches(0.7)
    add_icon_circle(slide, x + (step_w - circle_size) // 2, y_top, circle_size, step['color'], step['num'], font_size=22)

    # Title
    add_text(slide, x, y_top + Inches(0.9), step_w, Inches(0.4),
             step['title'], font_size=16, color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)

    # Description
    add_multiline_text(slide, x, y_top + Inches(1.4), step_w, Inches(1.0),
                       step['desc'].split('\n'), font_size=11, color=DARK_GRAY, line_spacing=1.3)
    for p in slide.shapes[-1].text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER

    # Arrow between steps
    if i < len(steps) - 1:
        add_arrow_right(slide, x + step_w + Inches(0.05), y_top + Inches(0.2), Inches(0.4), Inches(0.3), MID_GRAY)


# Status flow at bottom
status_y = Inches(5.0)
add_text(slide, Inches(0.8), status_y, Inches(10), Inches(0.4),
         'ステータスの遷移', font_size=16, color=BLACK, bold=True)

statuses = [
    ('要対応', RED, RED_LIGHT),
    ('対応中', AMBER, AMBER_LIGHT),
    ('完了', BLUE, BLUE_LIGHT),
]

sx = Inches(0.8)
for i, (label, color, bg) in enumerate(statuses):
    add_badge(slide, sx, status_y + Inches(0.6), label, bg, color, Inches(1.5), Inches(0.45), font_size=13)
    if i < len(statuses) - 1:
        add_arrow_right(slide, sx + Inches(1.6), status_y + Inches(0.68), Inches(0.5), Inches(0.28), MID_GRAY)
    sx += Inches(2.2)

# Cancelled (separate)
add_text(slide, sx + Inches(0.5), status_y + Inches(0.55), Inches(3), Inches(0.5),
         '※ キャンセルはいつでも可能', font_size=11, color=MID_GRAY)
add_badge(slide, sx + Inches(0.5), status_y + Inches(1.1), 'キャンセル', LIGHT_GRAY, DARK_GRAY, Inches(1.5), Inches(0.4), font_size=12)


# =====================================================
# Slide 4: Step 1 - Task Confirmation
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), fill_color=INDIGO)

add_icon_circle(slide, Inches(0.8), Inches(0.35), Inches(0.5), INDIGO, '1', font_size=18)
add_text(slide, Inches(1.5), Inches(0.35), Inches(8), Inches(0.5),
         'STEP 1: 当日タスクの確認', font_size=28, color=BLACK, bold=True)

# Left panel - instructions
left_x = Inches(0.8)
left_w = Inches(5.8)

add_rect(slide, left_x, Inches(1.3), left_w, Inches(5.5), fill_color=VERY_LIGHT, corner_radius=True)

instructions = [
    ('PMS管理画面にログイン', 'ブラウザで pms.tiramis.co.jp にアクセスし、\n自分のアカウントでログインします。'),
    ('サイドバーから「中継/回収管理」を選択', '左側メニューの OPERATIONS グループにある\n「中継/回収管理」をクリックします。'),
    ('当日の日付を確認', '画面上部のフィルタバーで対象日が\n本日の日付になっていることを確認します。'),
    ('自分のタスクを絞り込む', 'ドライバーフィルタで自分の名前を選択し、\n自分に割り当てられたタスクのみ表示します。'),
]

for i, (title, desc) in enumerate(instructions):
    iy = Inches(1.6) + Inches(1.25) * i
    # Number badge
    add_icon_circle(slide, left_x + Inches(0.3), iy, Inches(0.35), INDIGO, str(i + 1), font_size=12)
    add_text(slide, left_x + Inches(0.8), iy - Inches(0.02), Inches(4.5), Inches(0.35),
             title, font_size=13, color=BLACK, bold=True)
    add_multiline_text(slide, left_x + Inches(0.8), iy + Inches(0.3), Inches(4.5), Inches(0.8),
                       desc.split('\n'), font_size=10, color=DARK_GRAY, line_spacing=1.3)

# Right panel - what to check
right_x = Inches(7.0)
right_w = Inches(5.5)

add_text(slide, right_x, Inches(1.3), right_w, Inches(0.4),
         '確認すべき項目', font_size=16, color=BLACK, bold=True)

check_items = [
    ('種別', '中継 / 回収 / 全中継', ORANGE),
    ('到着時間枠', '配布員と合流する時間帯', BLUE),
    ('場所', '受け渡し場所（住所 or 地図）', GREEN),
    ('配布員', '対象の配布員名・エリア', PURPLE),
    ('チラシ情報', '必要なチラシの種類・部数', INDIGO),
    ('メモ', '特記事項（駐車場所等）', DARK_GRAY),
]

for i, (label, desc, color) in enumerate(check_items):
    cy = Inches(1.9) + Inches(0.75) * i
    add_rect(slide, right_x, cy, right_w, Inches(0.65), fill_color=VERY_LIGHT, border_color=RGBColor(0xE2, 0xE8, 0xF0), border_width=Pt(1), corner_radius=True)
    # Color dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, right_x + Inches(0.15), cy + Inches(0.2), Inches(0.25), Inches(0.25))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    add_text(slide, right_x + Inches(0.5), cy + Inches(0.05), Inches(1.5), Inches(0.35),
             label, font_size=12, color=BLACK, bold=True)
    add_text(slide, right_x + Inches(0.5), cy + Inches(0.35), Inches(4.5), Inches(0.3),
             desc, font_size=10, color=DARK_GRAY)


# =====================================================
# Slide 5: Step 2 - Route Optimization
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), fill_color=INDIGO)

add_icon_circle(slide, Inches(0.8), Inches(0.35), Inches(0.5), BLUE, '2', font_size=18)
add_text(slide, Inches(1.5), Inches(0.35), Inches(8), Inches(0.5),
         'STEP 2: ルート最適化', font_size=28, color=BLACK, bold=True)

add_text(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.4),
         '複数タスクを効率的に巡回するための最適ルートを自動計算します', font_size=14, color=MID_GRAY)

# Left: How to use
add_rect(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), fill_color=VERY_LIGHT, corner_radius=True)
add_text(slide, Inches(1.1), Inches(2.0), Inches(5), Inches(0.4),
         '使い方', font_size=16, color=BLACK, bold=True)

route_steps = [
    '「ルート最適化」ボタンをクリック',
    '優先順位を選択（回収優先 / 中継優先 / 時間優先）',
    '出発地点を選択（高田馬場 or 現在地）',
    'マップ上でルートと巡回順を確認',
    '「順番を適用」で優先順位を保存',
    '「ナビを開始」でGoogle Mapsナビへ',
]

for i, step_text in enumerate(route_steps):
    sy = Inches(2.6) + Inches(0.6) * i
    add_icon_circle(slide, Inches(1.1), sy, Inches(0.3), BLUE, str(i + 1), font_size=10)
    add_text(slide, Inches(1.6), sy - Inches(0.02), Inches(4.5), Inches(0.35),
             step_text, font_size=12, color=DARK_GRAY)

# Right: Priority modes
right_x = Inches(6.8)
add_text(slide, right_x, Inches(1.8), Inches(5), Inches(0.4),
         '優先モード', font_size=16, color=BLACK, bold=True)

modes = [
    ('回収優先', '回収タスクを先に巡回し、\nその後に中継タスクを巡回', PURPLE, PURPLE_LIGHT),
    ('中継優先', '中継タスクを先に巡回し、\nその後に回収タスクを巡回', ORANGE, ORANGE_LIGHT),
    ('時間優先', '到着時間枠を考慮して\n最短ルートを自動計算', INDIGO, INDIGO_LIGHT),
]

for i, (title, desc, color, bg) in enumerate(modes):
    my = Inches(2.4) + Inches(1.3) * i
    add_rect(slide, right_x, my, Inches(5.5), Inches(1.1), fill_color=bg, corner_radius=True)
    add_badge(slide, right_x + Inches(0.2), my + Inches(0.15), title, color, WHITE, Inches(1.4), Inches(0.35), font_size=11)
    add_multiline_text(slide, right_x + Inches(1.8), my + Inches(0.1), Inches(3.5), Inches(0.9),
                       desc.split('\n'), font_size=11, color=DARK_GRAY, line_spacing=1.3)

# Bottom note
add_rect(slide, right_x, Inches(6.3), Inches(5.5), Inches(0.5), fill_color=AMBER_LIGHT, corner_radius=True)
add_text(slide, right_x + Inches(0.2), Inches(6.35), Inches(5), Inches(0.4),
         '⚠ 時間指定ありのタスクは自動的にその時間枠に固定されます', font_size=10, color=AMBER)


# =====================================================
# Slide 6: Step 3 & 4 - Navigation & Handoff
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), fill_color=INDIGO)

# Split into two halves
# Left: Step 3
add_icon_circle(slide, Inches(0.8), Inches(0.35), Inches(0.5), GREEN, '3', font_size=18)
add_text(slide, Inches(1.5), Inches(0.35), Inches(4), Inches(0.5),
         'STEP 3: ナビゲーション', font_size=24, color=BLACK, bold=True)

add_rect(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(5.6), fill_color=GREEN_LIGHT, corner_radius=True)

nav_items = [
    ('ナビ開始ボタン', 'ルート最適化画面の「ナビを開始」ボタンを\nタップすると、Google Maps が起動します。'),
    ('ウェイポイント自動設定', '全タスクの場所が経由地として自動設定され、\n最適化された順序でナビゲーションします。'),
    ('到着時間の確認', '各地点の到着予想時刻（ETA）が\nルート画面に表示されています。'),
    ('時間指定タスクに注意', '時間枠が設定されているタスクは\nその時間に合わせて到着してください。'),
]

for i, (title, desc) in enumerate(nav_items):
    ny = Inches(1.5) + Inches(1.3) * i
    add_rect(slide, Inches(1.0), ny, Inches(5.1), Inches(1.1), fill_color=WHITE, corner_radius=True)
    add_text(slide, Inches(1.3), ny + Inches(0.1), Inches(4.5), Inches(0.3),
             title, font_size=12, color=GREEN, bold=True)
    add_multiline_text(slide, Inches(1.3), ny + Inches(0.4), Inches(4.5), Inches(0.6),
                       desc.split('\n'), font_size=10, color=DARK_GRAY, line_spacing=1.3)

# Right: Step 4
add_icon_circle(slide, Inches(7.0), Inches(0.35), Inches(0.5), ORANGE, '4', font_size=18)
add_text(slide, Inches(7.7), Inches(0.35), Inches(4), Inches(0.5),
         'STEP 4: 中継/回収の実施', font_size=24, color=BLACK, bold=True)

add_rect(slide, Inches(7.0), Inches(1.2), Inches(5.5), Inches(5.6), fill_color=ORANGE_LIGHT, corner_radius=True)

handoff_items = [
    ('中継の場合', '・配布員と合流し、チラシを手渡す\n・チラシの種類と部数を確認\n・受け渡し場所で待ち合わせ'),
    ('回収の場合', '・配布員から余ったチラシを受け取る\n・チラシの種類と枚数を記録\n・車両に積み込み'),
    ('全中継の場合', '・全チラシをまとめて配布員に渡す\n・チラシリスト全体を照合\n・不足がないか確認'),
    ('ステータス変更', '受け渡し完了後、\nPMS画面でステータスを\n「対応中」→「完了」に変更'),
]

for i, (title, desc) in enumerate(handoff_items):
    hy = Inches(1.5) + Inches(1.3) * i
    add_rect(slide, Inches(7.2), hy, Inches(5.1), Inches(1.1), fill_color=WHITE, corner_radius=True)
    add_text(slide, Inches(7.5), hy + Inches(0.1), Inches(4.5), Inches(0.3),
             title, font_size=12, color=ORANGE, bold=True)
    add_multiline_text(slide, Inches(7.5), hy + Inches(0.35), Inches(4.5), Inches(0.7),
                       desc.split('\n'), font_size=10, color=DARK_GRAY, line_spacing=1.2)


# =====================================================
# Slide 7: Step 5 - Evidence Photos
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), fill_color=INDIGO)

add_icon_circle(slide, Inches(0.8), Inches(0.35), Inches(0.5), PURPLE, '5', font_size=18)
add_text(slide, Inches(1.5), Inches(0.35), Inches(8), Inches(0.5),
         'STEP 5: エビデンス写真の撮影・アップロード', font_size=28, color=BLACK, bold=True)

add_text(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.4),
         '中継/回収の実施証拠として、現場の写真を撮影しPMSにアップロードします', font_size=14, color=MID_GRAY)

# Left: Why evidence?
add_rect(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.5), fill_color=PURPLE_LIGHT, corner_radius=True)
add_text(slide, Inches(1.1), Inches(2.0), Inches(5), Inches(0.4),
         'なぜエビデンスが必要？', font_size=16, color=PURPLE, bold=True)

why_items = [
    '業務の実施記録として保存',
    'チラシの受け渡し確認の証拠',
    '問題発生時のトラブルシューティング',
    '業務品質の管理・改善',
]
for i, item in enumerate(why_items):
    wy = Inches(2.6) + Inches(0.4) * i
    add_text(slide, Inches(1.3), wy, Inches(5), Inches(0.35),
             f'✓  {item}', font_size=12, color=DARK_GRAY)

# Right: How to upload
add_rect(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(2.5), fill_color=INDIGO_LIGHT, corner_radius=True)
add_text(slide, Inches(7.3), Inches(2.0), Inches(5), Inches(0.4),
         'アップロード手順', font_size=16, color=INDIGO, bold=True)

upload_steps_text = [
    '① タスクをタップして編集画面を開く',
    '② 「エビデンス写真」セクションまでスクロール',
    '③ 「写真を追加」ボタンをタップ',
    '④ カメラで撮影 or ギャラリーから選択',
    '⑤ 複数枚同時にアップロード可能',
]
for i, item in enumerate(upload_steps_text):
    uy = Inches(2.6) + Inches(0.4) * i
    add_text(slide, Inches(7.3), uy, Inches(5), Inches(0.35),
             item, font_size=12, color=DARK_GRAY)

# Bottom: Photo tips
add_text(slide, Inches(0.8), Inches(4.6), Inches(12), Inches(0.4),
         '撮影のポイント', font_size=16, color=BLACK, bold=True)

tips_data = [
    ('受け渡し場所', '合流した場所の全景を\n撮影してください', '📍'),
    ('チラシの確認', 'チラシの種類・束が\n分かるように撮影', '📄'),
    ('車両とチラシ', '積み込み/積み降ろし時の\n車両とチラシを撮影', '🚗'),
    ('明るく鮮明に', '暗い場所ではフラッシュを\n使用してください', '💡'),
]

for i, (title, desc, icon) in enumerate(tips_data):
    tx = Inches(0.8) + Inches(3.1) * i
    ty = Inches(5.1)
    add_rect(slide, tx, ty, Inches(2.8), Inches(1.8), fill_color=VERY_LIGHT, border_color=RGBColor(0xE2, 0xE8, 0xF0), border_width=Pt(1), corner_radius=True)
    add_text(slide, tx + Inches(0.2), ty + Inches(0.15), Inches(0.5), Inches(0.5), icon, font_size=24)
    add_text(slide, tx + Inches(0.2), ty + Inches(0.65), Inches(2.4), Inches(0.35),
             title, font_size=12, color=BLACK, bold=True)
    add_multiline_text(slide, tx + Inches(0.2), ty + Inches(1.0), Inches(2.4), Inches(0.7),
                       desc.split('\n'), font_size=10, color=DARK_GRAY, line_spacing=1.3)


# =====================================================
# Slide 8: Step 6 - Completion & Carryover
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), fill_color=INDIGO)

add_icon_circle(slide, Inches(0.8), Inches(0.35), Inches(0.5), INDIGO_DARK, '6', font_size=18)
add_text(slide, Inches(1.5), Inches(0.35), Inches(8), Inches(0.5),
         'STEP 6: 完了報告 & 翌日繰越', font_size=28, color=BLACK, bold=True)

# Left: Completion
add_rect(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(5.5), fill_color=BLUE_LIGHT, corner_radius=True)
add_text(slide, Inches(1.1), Inches(1.5), Inches(5), Inches(0.4),
         '完了報告の手順', font_size=18, color=BLUE, bold=True)

completion_steps = [
    ('ステータス変更', '各タスクのステータスドロップダウンから\n「完了」を選択します。'),
    ('エビデンス確認', 'エビデンス写真がアップロード\nされていることを確認します。'),
    ('メモの記入', '特筆事項があれば、メモ欄に\n記録を残してください。'),
    ('全タスク確認', '当日の全タスクが「完了」または\n「キャンセル」になっていることを確認。'),
]

for i, (title, desc) in enumerate(completion_steps):
    cy = Inches(2.2) + Inches(1.1) * i
    add_icon_circle(slide, Inches(1.1), cy + Inches(0.05), Inches(0.35), BLUE, str(i + 1), font_size=12)
    add_text(slide, Inches(1.6), cy, Inches(4.5), Inches(0.35),
             title, font_size=13, color=BLACK, bold=True)
    add_multiline_text(slide, Inches(1.6), cy + Inches(0.35), Inches(4.5), Inches(0.6),
                       desc.split('\n'), font_size=10, color=DARK_GRAY, line_spacing=1.3)

# Right: Carryover
add_rect(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(3.0), fill_color=AMBER_LIGHT, corner_radius=True)
add_text(slide, Inches(7.3), Inches(1.5), Inches(5), Inches(0.4),
         '翌日繰越について', font_size=18, color=AMBER, bold=True)

add_multiline_text(slide, Inches(7.3), Inches(2.1), Inches(5), Inches(2.0), [
    '当日中に完了できなかったタスクは、',
    '「翌日に繰越」ボタンで翌日に移動できます。',
    '',
    '【操作手順】',
    '1. 対象タスクの「翌日に繰越」ボタンをクリック',
    '2. 確認ダイアログで「繰越」を選択',
    '3. タスクが翌日の一覧に移動します',
], font_size=11, color=DARK_GRAY, line_spacing=1.3)

# Bottom right: Important notes
add_rect(slide, Inches(7.0), Inches(4.6), Inches(5.5), Inches(2.2), fill_color=RED_LIGHT, corner_radius=True)
add_text(slide, Inches(7.3), Inches(4.8), Inches(5), Inches(0.4),
         '注意事項', font_size=16, color=RED, bold=True)

notes = [
    '・ステータス変更は速やかに行ってください',
    '・エビデンス写真は必ずアップロードしてください',
    '・問題が発生した場合はメモに記録してください',
    '・緊急時は管理者に連絡してください',
]
for i, note in enumerate(notes):
    add_text(slide, Inches(7.3), Inches(5.3) + Inches(0.35) * i, Inches(5), Inches(0.35),
             note, font_size=11, color=RED)


# =====================================================
# Slide 9: Tips & FAQ
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), fill_color=INDIGO)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
         'よくある質問 & Tips', font_size=30, color=BLACK, bold=True)

faqs = [
    ('Q. 配布員が来ない場合は？', 'メモに状況を記録し、管理者に連絡してください。\nステータスは「対応中」のままにしておきます。'),
    ('Q. 場所が分からない場合は？', 'タスクの地図アイコンをタップすると\nGoogle Maps で場所を確認できます。'),
    ('Q. タスクが追加/変更された場合は？', '画面を更新（プルダウン更新 or リロード）すると\n最新のタスク一覧が表示されます。'),
    ('Q. 写真のアップロードに失敗する場合は？', '電波の良い場所で再試行してください。\nWi-Fi環境での再アップロードも有効です。'),
    ('Q. 2つのタスクが同じ時間帯の場合は？', 'ルート最適化で最適な巡回順を確認してください。\n必要に応じて管理者に調整を依頼します。'),
]

for i, (q, a) in enumerate(faqs):
    fy = Inches(1.3) + Inches(1.15) * i
    fw = Inches(11.7)
    add_rect(slide, Inches(0.8), fy, fw, Inches(1.0), fill_color=VERY_LIGHT, border_color=RGBColor(0xE2, 0xE8, 0xF0), border_width=Pt(1), corner_radius=True)
    add_text(slide, Inches(1.1), fy + Inches(0.08), Inches(4), Inches(0.35),
             q, font_size=12, color=INDIGO, bold=True)
    add_multiline_text(slide, Inches(5.5), fy + Inches(0.08), Inches(6.5), Inches(0.8),
                       a.split('\n'), font_size=11, color=DARK_GRAY, line_spacing=1.3)


# =====================================================
# Slide 10: End
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, VERY_LIGHT)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), fill_color=INDIGO)

# Center card
card_w = Inches(8)
card_h = Inches(3.5)
card_left = (W - card_w) // 2
card_top = Inches(2.0)
add_rect(slide, card_left, card_top, card_w, card_h, fill_color=WHITE, border_color=RGBColor(0xE2, 0xE8, 0xF0), border_width=Pt(1), corner_radius=True)

add_text(slide, card_left, card_top + Inches(0.8), card_w, Inches(0.8),
         'お疲れさまでした', font_size=36, color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)

add_text(slide, card_left, card_top + Inches(1.8), card_w, Inches(0.5),
         '不明点や改善要望は管理者までご連絡ください', font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_text(slide, card_left, card_top + Inches(2.6), card_w, Inches(0.5),
         '株式会社ティラミス  |  PMS Pro', font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# Save
output_dir = os.path.expanduser('~/Downloads/PMS_設計書')
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, '中継回収_担当者マニュアル.pptx')
prs.save(output_path)
print(f'Saved: {output_path}')
